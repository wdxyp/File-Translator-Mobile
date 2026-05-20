import warnings
# 静默 openpyxl 关于 DrawingML 支持不全的警告
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl.reader.drawings")

import sys
import argparse
try:
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk
    TK_AVAILABLE = True
except Exception:
    tk = None
    filedialog = None
    messagebox = None
    ttk = None
    TK_AVAILABLE = False
from pptx import Presentation
try:
    import pandas as pd
except Exception:
    pd = None
import os
import hashlib
import random
import string
import requests
from openpyxl import load_workbook, Workbook
from openpyxl.styles import Font, Border, Alignment, PatternFill
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.oxml.ns import nsmap as DOCX_NSMAP
from datetime import datetime
import re
import time
import threading
import shutil

# 请替换为你自己的 APP ID 和密钥（建议通过环境变量注入，避免提交到仓库）
APP_ID = os.getenv('BAIDU_APP_ID') or 'YOUR_APP_ID'
SECRET_KEY = os.getenv('BAIDU_SECRET_KEY') or 'YOUR_SECRET_KEY'
API_URL = 'https://api.fanyi.baidu.com/api/trans/vip/translate'

def configure_baidu(app_id=None, secret_key=None):
    global APP_ID, SECRET_KEY
    if app_id:
        APP_ID = str(app_id).strip()
    if secret_key:
        SECRET_KEY = str(secret_key).strip()

# ==========================================
# V2.11 (百度长语句翻译 - 稳定最终版)
# ==========================================

# 全局变量
revision_map = {}
original_texts = []
translated_texts = []

class _ValueBox:
    def __init__(self, value=None):
        self._value = value
    def get(self):
        return self._value
    def set(self, value):
        self._value = value

FROM_LANG_MAP = {
    'zh2ko': 'zh', 'ko2zh': 'kor', 'ko2en': 'kor', 'zh2en': 'zh', 'en2zh': 'en',
    'zh_tw2en': 'zh', 'en2zh_tw': 'en', 'zh2ja': 'zh', 'ja2zh': 'ja',
    'en2ko': 'en', 'vi2zh': 'vie', 'zh2vi': 'zh', 'ko2vi': 'kor'
}
TO_LANG_MAP = {
    'zh2ko': 'kor', 'ko2zh': 'zh', 'ko2en': 'en', 'zh2en': 'en', 'en2zh': 'zh',
    'zh_tw2en': 'en', 'en2zh_tw': 'zh', 'zh2ja': 'ja', 'ja2zh': 'zh',
    'en2ko': 'kor', 'vi2zh': 'zh', 'zh2vi': 'vie', 'ko2vi': 'vie'
}
TARGET_FONT_BY_TO_LANG = {'zh': '微软雅黑', 'kor': 'Malgun Gothic'}
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def get_target_to_lang():
    try:
        direction = translation_direction.get()
    except Exception:
        return None
    return TO_LANG_MAP.get(direction)

def get_target_font_name():
    to_lang = get_target_to_lang()
    if not to_lang:
        return None
    return TARGET_FONT_BY_TO_LANG.get(to_lang)

def xpath_with_ns(element, expr):
    if element is None:
        return []
    try:
        return element.xpath(expr, namespaces=DOCX_NSMAP)
    except TypeError:
        return element.xpath(expr)

def set_docx_r_element_font(r_element, font_name):
    if r_element is None or not font_name:
        return
    rPr = r_element.get_or_add_rPr()
    rFonts = rPr.rFonts
    if rFonts is None:
        rFonts = OxmlElement('w:rFonts')
        rPr.append(rFonts)
    rFonts.set(qn('w:ascii'), font_name)
    rFonts.set(qn('w:hAnsi'), font_name)
    rFonts.set(qn('w:eastAsia'), font_name)

def set_docx_run_font(run, font_name):
    if not run or not font_name:
        return
    run.font.name = font_name
    set_docx_r_element_font(run._r, font_name)

def set_drawingml_rpr_element_font(a_rPr_element, font_name):
    if a_rPr_element is None or not font_name:
        return
    for local_name, tag in (('latin', 'a:latin'), ('ea', 'a:ea'), ('cs', 'a:cs')):
        el = a_rPr_element.find(f'{{{A_NS}}}{local_name}')
        if el is None:
            el = OxmlElement(tag)
            a_rPr_element.append(el)
        el.set('typeface', font_name)

def set_drawingml_r_element_font(a_r_element, font_name):
    if a_r_element is None or not font_name:
        return
    a_rPr = a_r_element.find(f'{{{A_NS}}}rPr')
    if a_rPr is None:
        a_rPr = OxmlElement('a:rPr')
        a_r_element.insert(0, a_rPr)
    set_drawingml_rpr_element_font(a_rPr, font_name)

def set_pptx_run_font(run, font_name):
    if not run or not font_name:
        return
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        r = run._r
        if hasattr(r, 'get_or_add_rPr'):
            rPr = r.get_or_add_rPr()
        else:
            rPr = r.find(f'{{{A_NS}}}rPr')
            if rPr is None:
                rPr = OxmlElement('a:rPr')
                r.insert(0, rPr)
        set_drawingml_rpr_element_font(rPr, font_name)
    except Exception:
        return

def set_pptx_paragraph_default_font(paragraph, font_name):
    if not paragraph or not font_name:
        return
    try:
        p = paragraph._p
        pPr = p.find(f'{{{A_NS}}}pPr')
        if pPr is None:
            pPr = OxmlElement('a:pPr')
            p.insert(0, pPr)
        defRPr = pPr.find(f'{{{A_NS}}}defRPr')
        if defRPr is None:
            defRPr = OxmlElement('a:defRPr')
            pPr.append(defRPr)
        set_drawingml_rpr_element_font(defRPr, font_name)
    except Exception:
        return

def load_revision_dict(file_path="revision.md", silent=False):
    """
    自动加载校准词典 (revision.md)
    """
    mapping = {}
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line or line.startswith("#"): continue
                    if not re.match(r'^[-*+]\s*', line): continue
                    content = re.sub(r'^[-*+]\s*', '', line)
                    if "格式" in content: continue
                    
                    if ":" in content:
                        parts = content.split(":", 1)
                        if len(parts) == 2:
                            err, corr = parts
                            if err.strip(): mapping[err.strip()] = corr.strip()
                    elif "：" in content:
                        parts = content.split("：", 1)
                        if len(parts) == 2:
                            err, corr = parts
                            if err.strip(): mapping[err.strip()] = corr.strip()
            if not silent:
                print(f"[系统] 自动加载校准文件成功，共 {len(mapping)} 条有效规则。")
        except Exception as e:
            print(f"[错误] 加载校准文件失败: {e}")
    return mapping

def apply_revisions(text):
    """
    在翻译结果上应用校准映射
    """
    if not text or not revision_map: return text
    sorted_keys = sorted(revision_map.keys(), key=len, reverse=True)
    result_text = text
    for err in sorted_keys:
        if err in result_text:
            result_text = result_text.replace(err, revision_map[err])
    return result_text

def baidu_translate(q, from_lang, to_lang):
    if not q or not q.strip(): return q
    max_retries = 3
    retry_delay = 2
    retries = 0
    while retries < max_retries:
        try:
            salt = ''.join(random.choices(string.ascii_letters + string.digits, k=16))
            sign_str = APP_ID + q + salt + SECRET_KEY
            sign = hashlib.md5(sign_str.encode()).hexdigest()
            params = {
                'q': q, 'from': from_lang, 'to': to_lang,
                'appid': APP_ID, 'salt': salt, 'sign': sign
            }
            response = requests.get(API_URL, params=params, timeout=15)
            response.encoding = 'utf-8'
            result = response.json()
            if 'trans_result' in result:
                return result['trans_result'][0]['dst']
            else:
                return q
        except Exception:
            retries += 1
            if retries < max_retries:
                time.sleep(retry_delay)
    return q

def get_translation(text):
    if not text or not text.strip(): return text
    direction = translation_direction.get()
    from_lang = FROM_LANG_MAP.get(direction, 'auto')
    to_lang = TO_LANG_MAP.get(direction, 'auto')
    
    # 百度长句翻译
    translated_text = baidu_translate(text, from_lang, to_lang)
    # 事后校准
    translated_text = apply_revisions(translated_text)
    return translated_text

def _ppt_normalize_linebreaks(text):
    if not text:
        return text
    return str(text).replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")

def _ppt_denormalize_linebreaks(text):
    if not text:
        return text
    return str(text).replace("\n", "\v")

def translate_ppt_paragraph(paragraph):
    full_text = paragraph.text
    if not full_text or not str(full_text).strip():
        return

    normalized = _ppt_normalize_linebreaks(full_text)
    translated = get_translation(normalized)

    original_texts.append(normalized)
    translated_texts.append(translated)

    if append_translation.get():
        result_norm = append_translation_to_original(normalized, translated)
    else:
        result_norm = translated

    result_text = _ppt_denormalize_linebreaks(result_norm)

    if paragraph.runs:
        first_run = paragraph.runs[0]
    else:
        first_run = paragraph.add_run()

    original_font_name = first_run.font.name
    original_size = first_run.font.size
    original_bold = first_run.font.bold
    original_italic = first_run.font.italic
    original_underline = first_run.font.underline

    first_run.text = result_text
    for r in paragraph.runs[1:]:
        r.text = ""

    target_font_name = get_target_font_name()
    if target_font_name:
        set_pptx_run_font(first_run, target_font_name)
        set_pptx_paragraph_default_font(paragraph, target_font_name)
    elif original_font_name:
        first_run.font.name = original_font_name
        first_run.font.size = original_size
        first_run.font.bold = original_bold
        first_run.font.italic = original_italic
        first_run.font.underline = original_underline

def append_translation_to_original(text, translated_text, cell=None):
    text = text.strip()
    translated_text = translated_text.strip()
    result = f"{text}\n{translated_text}" if text and translated_text else (text or translated_text)
    if cell:
        cell.alignment = Alignment(wrap_text=True, vertical='center')
        ws = cell.parent
        row_num = cell.row
        original_height = ws.row_dimensions[row_num].height
        ws.row_dimensions[row_num].height = (original_height * 2) if original_height else 30
    return result

# ==========================================
# 文件处理逻辑 (完全基于 V2.9 的稳定代码)
# ==========================================

def translate_shape_for_ppt(shape):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            translate_ppt_paragraph(paragraph)
    elif shape.shape_type == 6:  # 组合形状
        for sub_shape in shape.shapes:
            translate_shape_for_ppt(sub_shape)
    elif shape.has_table:  # 处理表格形状
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    translate_ppt_paragraph(paragraph)

def update_ui_status(msg):
    """线程安全地更新 UI 状态"""
    r = globals().get("root")
    lbl = globals().get("status_label")
    if r is not None and hasattr(r, "after") and lbl is not None:
        r.after(0, lambda: lbl.config(text=msg))
    else:
        print(f"[进度] {msg}")

def translate_ppt(input_file, output_file):
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        msg = f"正在翻译 PPT: 第 {i}/{total_slides} 页..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        for shape in slide.shapes:
            translate_shape_for_ppt(shape)
    
    print(f"[系统] 正在保存 PPT 文件...")
    prs.save(output_file)
    save_to_corpus(original_texts, translated_texts)
    print(f"[完成] 保存修改后的 PPT 文件完成!")

def clean_sheet_name(name):
    if not name: return "Sheet"
    invalid_chars = r'[\\/?:*\[\](){}<>|"\']'
    clean_name = re.sub(invalid_chars, '', name)
    return clean_name[:31]

def translate_excel_xlsx(input_file, output_file):
    """处理.xlsx格式的Excel文件"""
    # 加载工作簿 (keep_vba=True 尝试保留宏和部分绘图元数据)
    wb = load_workbook(input_file, keep_vba=True)
    # 收集所有已存在的工作表名称（包括原始和已处理的）
    existing_sheet_names = set()
    total_sheets = len(wb.sheetnames)
    for i, sheet_name in enumerate(wb.sheetnames, 1):
        msg = f"正在翻译 Excel: 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        
        # 翻译工作表名称 (还原 V2.9 逻辑并增加翻译)
        translated_sheet_name = get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        
        # 确保工作表名称唯一
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)
        
        ws = wb[sheet_name]
        # 修改工作表名称
        ws.title = unique_sheet_name
        
        # 获取总行数用于进度（大致）
        total_rows = ws.max_row
        for row_idx, row in enumerate(ws.iter_rows(), 1):
            if row_idx % 10 == 0:
                row_msg = f"处理中: {unique_sheet_name} 第 {row_idx}/{total_rows} 行..."
                print(f"[进度] {row_msg}")
                update_ui_status(row_msg)
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    original_value = str(cell.value)
                    normalized_value = original_value.replace("\r\n", "\n").replace("\r", "\n")
                    translated_text = get_translation(normalized_value)

                    original_texts.append(normalized_value)
                    translated_texts.append(translated_text)
                    
                    # 严格还原 V2.9 样式保护逻辑
                    # 保存原格式
                    original_font = Font(**cell.font.__dict__)
                    original_border = Border(**cell.border.__dict__)
                    original_alignment = Alignment(**cell.alignment.__dict__)
                    original_fill = PatternFill(**cell.fill.__dict__)
                    
                    if append_translation.get():
                        cell.value = append_translation_to_original(normalized_value, translated_text, cell)
                    else:
                        cell.value = translated_text
                    
                    # 恢复原格式
                    target_font_name = get_target_font_name()
                    if target_font_name:
                        try:
                            cell.font = original_font.copy(name=target_font_name)
                        except Exception:
                            cell.font = Font(
                                name=target_font_name,
                                size=original_font.size,
                                bold=original_font.bold,
                                italic=original_font.italic,
                                underline=original_font.underline,
                                color=original_font.color,
                            )
                    else:
                        cell.font = original_font
                    cell.border = original_border
                    cell.alignment = original_alignment
                    cell.fill = original_fill
                    # 设置单元格自动换行
                    cell.alignment = Alignment(wrap_text=True, vertical='center')
    # 保存修改后的工作簿
    wb.save(output_file)
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)

def translate_excel_xls(input_file, output_file):
    """处理.xls格式的Excel文件，使用pandas库，保存为.xlsx格式"""
    if pd is None:
        raise RuntimeError("处理 .xls 需要 pandas（以及其依赖）。请改用 .xlsx，或在运行环境中安装 pandas。")
    base_output = os.path.splitext(output_file)[0]
    output_file_xlsx = base_output + '.xlsx'
    target_font_name = get_target_font_name()
    # 读取所有工作表
    excel_file = pd.ExcelFile(input_file)
    # 创建一个新的ExcelWriter对象，使用默认引擎
    writer = pd.ExcelWriter(output_file_xlsx, engine='openpyxl')
    # 收集所有已存在的工作表名称（包括原始和已处理的）
    existing_sheet_names = set()
    total_sheets = len(excel_file.sheet_names)
    for i, sheet_name in enumerate(excel_file.sheet_names, 1):
        msg = f"正在翻译 Excel(.xls): 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        
        # 翻译工作表名称
        translated_sheet_name = get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        
        # 确保工作表名称唯一
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)
        
        # 读取工作表
        df = pd.read_excel(input_file, sheet_name=sheet_name)
        # 遍历所有单元格进行翻译
        total_rows = len(df.index)
        for idx_num, idx in enumerate(df.index, 1):
            if idx_num % 10 == 0:
                row_msg = f"处理中: {unique_sheet_name} 第 {idx_num}/{total_rows} 行..."
                print(f"[进度] {row_msg}")
                update_ui_status(row_msg)
            for col in df.columns:
                cell_value = df.at[idx, col]
                if pd.notna(cell_value) and isinstance(cell_value, str):
                    original_value = str(cell_value)
                    normalized_value = original_value.replace("\r\n", "\n").replace("\r", "\n")
                    translated_text = get_translation(normalized_value)

                    original_texts.append(normalized_value)
                    translated_texts.append(translated_text)

                    df.at[idx, col] = append_translation_to_original(normalized_value, translated_text) if append_translation.get() else translated_text
        # 保存翻译后的工作表
        df.to_excel(writer, sheet_name=unique_sheet_name, index=False)
    # 保存工作簿
    writer.close()
    if target_font_name:
        try:
            wb = load_workbook(output_file_xlsx)
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            try:
                                cell.font = cell.font.copy(name=target_font_name)
                            except Exception:
                                cell.font = Font(
                                    name=target_font_name,
                                    size=cell.font.size,
                                    bold=cell.font.bold,
                                    italic=cell.font.italic,
                                    underline=cell.font.underline,
                                    color=cell.font.color,
                                )
            wb.save(output_file_xlsx)
        except Exception:
            pass
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)
    return output_file_xlsx

def translate_word(input_file, output_file):
    doc = Document(input_file)
    target_font_name = get_target_font_name()

    def translate_word_paragraph(paragraph):
        full_text = paragraph.text
        if not full_text or not str(full_text).strip():
            return

        normalized = str(full_text).replace("\r\n", "\n").replace("\r", "\n")
        t = get_translation(normalized)

        original_texts.append(normalized)
        translated_texts.append(t)

        result_text = append_translation_to_original(normalized, t) if append_translation.get() else t

        drawings = []
        for run in paragraph.runs:
            for inline in xpath_with_ns(run.element, './/w:drawing'):
                pic_elements = xpath_with_ns(inline, './/a:blip/@r:embed')
                if pic_elements:
                    drawings.append(inline)

        if paragraph.runs:
            first_run = paragraph.runs[0]
        else:
            first_run = paragraph.add_run()

        first_run.text = result_text
        if target_font_name:
            set_docx_run_font(first_run, target_font_name)

        for run in paragraph.runs[1:]:
            has_drawing = bool(xpath_with_ns(run.element, './/w:drawing'))
            if not has_drawing:
                run.text = ""

        for inline in drawings:
            new_run = paragraph.add_run()
            new_run._r.append(inline)
    
    # 处理普通段落 (完全还原 V2.9 逻辑)
    total_paragraphs = len(doc.paragraphs)
    for i, paragraph in enumerate(doc.paragraphs, 1):
        if i % 5 == 0:
            msg = f"正在翻译 Word: 段落 {i}/{total_paragraphs}..."
            print(f"[进度] {msg}")
            update_ui_status(msg)
        translate_word_paragraph(paragraph)
    
    # 处理表格 (完全还原 V2.9 逻辑)
    total_tables = len(doc.tables)
    for i, table in enumerate(doc.tables, 1):
        msg = f"正在翻译 Word: 表格 {i}/{total_tables}..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    translate_word_paragraph(paragraph)
    
    # 处理形状中的文字 (完全还原 V2.9 逻辑)
    for shape in doc.inline_shapes:
        try:
            # 通过 XML 访问形状中的文本
            if hasattr(shape, '_inline') and hasattr(shape._inline, 'graphic'):
                graphic_data = shape._inline.graphic.graphicData
                if hasattr(graphic_data, 'txBody') and graphic_data.txBody:
                    for p in graphic_data.txBody.p:
                        for r in p.r:
                            if hasattr(r, 't'):
                                text = r.t
                                if text and text.strip():  # 只翻译非空文本
                                    # 先获取翻译结果
                                    t = get_translation(text)
                                    
                                    # 收集翻译前后的文本
                                    original_texts.append(text)
                                    translated_texts.append(t)
                                    if append_translation.get():
                                        r.t = append_translation_to_original(text, t)
                                    else:
                                        r.t = t
                                    if target_font_name:
                                        set_drawingml_r_element_font(r, target_font_name)
        except Exception:
            continue
    doc_element_list = [doc.element]
    for section in doc.sections:
        doc_element_list.extend([
            section.header._element,
            section.footer._element,
            section.first_page_header._element,
            section.first_page_footer._element,
            section.even_page_header._element,
            section.even_page_footer._element,
        ])
    for doc_element in doc_element_list:
        for t_node in xpath_with_ns(doc_element, './/w:txbxContent//w:t'):
            text = t_node.text
            if not text or not text.strip():
                continue
            t = get_translation(text)
            original_texts.append(text)
            translated_texts.append(t)
            t_node.text = append_translation_to_original(text, t) if append_translation.get() else t
            if not target_font_name:
                continue
            r_element = t_node
            while r_element is not None and r_element.tag != qn('w:r'):
                r_element = r_element.getparent()
            if r_element is not None:
                set_docx_r_element_font(r_element, target_font_name)
        if not append_translation.get():
            for txbx in xpath_with_ns(doc_element, './/w:txbxContent'):
                t_nodes = xpath_with_ns(txbx, './/w:t')
                if not t_nodes:
                    continue
                combined = ''.join([(n.text or '') for n in t_nodes])
                revised = apply_revisions(combined)
                if revised != combined:
                    t_nodes[0].text = revised
                    for n in t_nodes[1:]:
                        n.text = ''
                    if target_font_name:
                        r_element = t_nodes[0]
                        while r_element is not None and r_element.tag != qn('w:r'):
                            r_element = r_element.getparent()
                        if r_element is not None:
                            set_docx_r_element_font(r_element, target_font_name)
        for t_node in doc_element.xpath(f'.//*[namespace-uri()="{A_NS}" and local-name()="txBody"]//*[namespace-uri()="{A_NS}" and local-name()="t"]'):
            text = t_node.text
            if not text or not text.strip():
                continue
            t = get_translation(text)
            original_texts.append(text)
            translated_texts.append(t)
            t_node.text = append_translation_to_original(text, t) if append_translation.get() else t
            if not target_font_name:
                continue
            a_r_element = t_node
            while a_r_element is not None and a_r_element.tag != f'{{{A_NS}}}r':
                a_r_element = a_r_element.getparent()
            if a_r_element is not None:
                set_drawingml_r_element_font(a_r_element, target_font_name)
        if not append_translation.get():
            for a_txbody in doc_element.xpath(f'.//*[namespace-uri()="{A_NS}" and local-name()="txBody"]'):
                t_nodes = a_txbody.xpath(f'.//*[namespace-uri()="{A_NS}" and local-name()="t"]')
                if not t_nodes:
                    continue
                combined = ''.join([(n.text or '') for n in t_nodes])
                revised = apply_revisions(combined)
                if revised != combined:
                    t_nodes[0].text = revised
                    for n in t_nodes[1:]:
                        n.text = ''
                    if target_font_name:
                        a_r_element = t_nodes[0]
                        while a_r_element is not None and a_r_element.tag != f'{{{A_NS}}}r':
                            a_r_element = a_r_element.getparent()
                        if a_r_element is not None:
                            set_drawingml_r_element_font(a_r_element, target_font_name)
        if target_font_name:
            for r_element in xpath_with_ns(doc_element, './/w:r'):
                set_docx_r_element_font(r_element, target_font_name)
            for r_element in xpath_with_ns(doc_element, './/w:txbxContent//w:r'):
                set_docx_r_element_font(r_element, target_font_name)
            for a_rpr_element in doc_element.xpath(f'.//*[namespace-uri()="{A_NS}" and (local-name()="rPr" or local-name()="defRPr" or local-name()="endParaRPr")]'):
                set_drawingml_rpr_element_font(a_rpr_element, target_font_name)

    doc.save(output_file)
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)

# --- 线程控制 ---

def start_translation():
    input_file = input_file_entry.get()
    output_folder = output_folder_entry.get()
    if not (input_file and output_folder):
        messagebox.showwarning("提示", "请完整选择输入文件和输出目录")
        return
    translate_button.config(state=tk.DISABLED, text="🚀 正在翻译，请稍候...")
    status_label.config(text="任务已启动，请查看终端进度...", foreground="#2980b9")
    thread = threading.Thread(target=run_translation_task, args=(input_file, output_folder))
    thread.daemon = True
    thread.start()

def run_translation_task(input_file, output_folder):
    try:
        global revision_map
        file_ext = os.path.splitext(input_file)[1].lower()
        custom_name = custom_filename_entry.get().strip()
        output_file = os.path.join(output_folder, (custom_name if custom_name else f"translated_v2.11_{os.path.basename(input_file).split('.')[0]}") + file_ext)
        
        # --- [V2.11 增强功能：物理克隆以保留图片和绘图] ---
        # 即使 V2.9 也不支持 Excel 绘图保留，这里通过物理复制尝试最大化兼容性
        try:
            if os.path.exists(output_file):
                os.remove(output_file) # 尝试删除旧文件，如果被占用会在这里报错
            shutil.copy2(input_file, output_file)
        except PermissionError:
            raise Exception(f"目标文件已被占用，请先关闭 Excel/Word/PPT: {os.path.basename(output_file)}")
        
        original_texts.clear()
        translated_texts.clear()
        revision_map = load_revision_dict("revision.md", silent=True)
        print(f"\n[系统] 开始翻译任务: {os.path.basename(input_file)}")
        print(f"[系统] 校准规则加载: {len(revision_map)} 条")
        
        # 统一使用 output_file 作为操作对象，实现“原地翻译”
        if file_ext in ['.ppt', '.pptx']: translate_ppt(output_file, output_file)
        elif file_ext == '.xlsx': translate_excel_xlsx(output_file, output_file)
        elif file_ext == '.xls': 
            # .xls 比较特殊，必须另存为 .xlsx
            output_file = translate_excel_xls(input_file, output_file)
        elif file_ext == '.docx': translate_word(output_file, output_file)
        
        print(f"[完成] 文件已保存至: {output_file}\n")
        root.after(0, lambda: translation_done_callback(output_file))
    except Exception as e:
        err_msg = str(e)
        print(f"[错误] 详情: {err_msg}")
        root.after(0, lambda: translation_failed_callback(err_msg))

def translation_done_callback(output_file):
    translate_button.config(state=tk.NORMAL, text="🚀 开始长句翻译任务")
    status_label.config(text=f"翻译任务已圆满完成！", foreground="#27ae60")
    messagebox.showinfo("成功", f"翻译完成！\n文件保存至：{output_file}")

def translation_failed_callback(error_msg):
    translate_button.config(state=tk.NORMAL, text="🚀 开始长句翻译任务")
    status_label.config(text=f"翻译过程出错", foreground="#e74c3c")
    messagebox.showerror("错误", f"发生异常：{error_msg}")

def save_to_corpus(orig, trans):
    if generate_corpus.get() and orig:
        if not os.path.exists('Corpus'): os.makedirs('Corpus')
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        try:
            direction = translation_direction.get()
        except Exception:
            direction = "unknown"
        corpus_file = f'Corpus/Corpus_v2.11_{direction}_{timestamp}.xlsx'

        to_lang = TO_LANG_MAP.get(direction)
        seen = set()
        filtered_orig = []
        filtered_trans = []
        for o, t in zip(orig, trans):
            if o is None or t is None:
                continue
            o = str(o).replace("\r\n", "\n").replace("\r", "\n").strip()
            t = str(t).replace("\r\n", "\n").replace("\r", "\n").strip()
            if not o or not t:
                continue
            if to_lang in ("zh", "zh_tw") and re.search(r"[\u4e00-\u9fff]", o):
                continue
            if o == t:
                continue
            if not re.search(r"[\uac00-\ud7a3]", o) and re.search(r"[A-Za-z0-9]", o):
                continue
            if not re.search(r"[\uac00-\ud7a3\u4e00-\u9fff]", o) and not re.search(r"[A-Za-z0-9]", o):
                continue
            key = o
            if key in seen:
                continue
            seen.add(key)
            filtered_orig.append(o)
            filtered_trans.append(t)

        if not filtered_orig:
            return

        if pd is not None:
            pd.DataFrame({'序号': range(1, len(filtered_orig)+1), '翻译前': filtered_orig, '翻译后': filtered_trans}).to_excel(corpus_file, index=False)
            return
        wb = Workbook()
        ws = wb.active
        ws.title = "Corpus"
        ws.append(["序号", "翻译前", "翻译后"])
        for i, (o, t) in enumerate(zip(filtered_orig, filtered_trans), 1):
            ws.append([i, o, t])
        wb.save(corpus_file)

def translate_file(input_file, output_dir=None, name="", direction="ko2zh", append=False, corpus=False, revision_file="revision.md", app_id=None, secret_key=None):
    global translation_direction, append_translation, generate_corpus, revision_map

    input_file = os.path.abspath(input_file)
    if not os.path.exists(input_file):
        raise FileNotFoundError(input_file)

    configure_baidu(app_id=app_id, secret_key=secret_key)

    output_dir = os.path.abspath(output_dir) if output_dir else os.path.dirname(input_file)
    os.makedirs(output_dir, exist_ok=True)

    translation_direction = _ValueBox(direction)
    append_translation = _ValueBox(bool(append))
    generate_corpus = _ValueBox(bool(corpus))

    original_texts.clear()
    translated_texts.clear()
    revision_map = load_revision_dict(revision_file, silent=True)

    file_ext = os.path.splitext(input_file)[1].lower()
    base_name = name.strip() if name and name.strip() else f"translated_{os.path.basename(input_file).split('.')[0]}"
    output_file = os.path.join(output_dir, base_name + file_ext)

    if os.path.exists(output_file):
        os.remove(output_file)
    shutil.copy2(input_file, output_file)

    if file_ext in ['.ppt', '.pptx']:
        translate_ppt(output_file, output_file)
    elif file_ext == '.xlsx':
        translate_excel_xlsx(output_file, output_file)
    elif file_ext == '.xls':
        output_file = translate_excel_xls(input_file, output_file)
    elif file_ext == '.docx':
        translate_word(output_file, output_file)
    else:
        raise ValueError(f"不支持的文件类型: {file_ext}")

    return output_file

def _run_cli():
    global translation_direction, append_translation, generate_corpus, revision_map
    parser = argparse.ArgumentParser(add_help=True)
    parser.add_argument("-i", "--input", required=True, help="输入文件路径（.pptx/.ppt/.docx/.xlsx/.xls）")
    parser.add_argument("-o", "--output-dir", default="", help="输出目录（默认同输入文件目录）")
    parser.add_argument("--name", default="", help="自定义输出文件名（不带扩展名）")
    parser.add_argument("--direction", default="ko2zh", choices=sorted(TO_LANG_MAP.keys()), help="翻译方向")
    parser.add_argument("--append", action="store_true", help="在原文下方追加翻译对照")
    parser.add_argument("--corpus", action="store_true", help="生成 Corpus 语料库文件")
    args = parser.parse_args()

    input_file = os.path.abspath(args.input)
    if not os.path.exists(input_file):
        raise SystemExit(f"输入文件不存在: {input_file}")

    output_dir = os.path.abspath(args.output_dir) if args.output_dir else os.path.dirname(input_file)
    os.makedirs(output_dir, exist_ok=True)

    translation_direction = _ValueBox(args.direction)
    append_translation = _ValueBox(bool(args.append))
    generate_corpus = _ValueBox(bool(args.corpus))

    original_texts.clear()
    translated_texts.clear()
    revision_map = load_revision_dict("revision.md", silent=True)
    print(f"\n[系统] 开始翻译任务: {os.path.basename(input_file)}")
    print(f"[系统] 校准规则加载: {len(revision_map)} 条")

    file_ext = os.path.splitext(input_file)[1].lower()
    base_name = args.name.strip() if args.name.strip() else f"translated_v2.11_{os.path.basename(input_file).split('.')[0]}"
    output_file = os.path.join(output_dir, base_name + file_ext)

    if os.path.exists(output_file):
        os.remove(output_file)
    shutil.copy2(input_file, output_file)

    if file_ext in ['.ppt', '.pptx']:
        translate_ppt(output_file, output_file)
    elif file_ext == '.xlsx':
        translate_excel_xlsx(output_file, output_file)
    elif file_ext == '.xls':
        output_file = translate_excel_xls(input_file, output_file)
    elif file_ext == '.docx':
        translate_word(output_file, output_file)
    else:
        raise SystemExit(f"不支持的文件类型: {file_ext}")

    print(f"[完成] 文件已保存至: {output_file}\n")

def _run_gui():
    global root, translation_direction, append_translation, generate_corpus
    global input_file_entry, output_folder_entry, custom_filename_entry
    global translate_button, status_label, revision_map
    root = tk.Tk()
    root.title("百度长语句翻译工具 V2.11 (完善版)")
    root.geometry("700x850")
    root.configure(bg="#f5f6fa")
    style = ttk.Style()
    style.theme_use('clam')
    style.configure("TFrame", background="#f5f6fa")
    style.configure("TLabel", background="#f5f6fa", font=("微软雅黑", 10))
    style.configure("Header.TLabel", font=("微软雅黑", 14, "bold"), foreground="#2f3640")
    main_frame = ttk.Frame(root, padding="20")
    main_frame.pack(fill="both", expand=True)
    ttk.Label(main_frame, text="百度长语句翻译 & 自动校准系统", style="Header.TLabel").pack(pady=(0, 20))
    file_card = ttk.LabelFrame(main_frame, text=" 文件设置 ", padding=15)
    file_card.pack(fill="x", pady=10)
    ttk.Label(file_card, text="待翻译文件:").grid(row=0, column=0, sticky="w", pady=5)
    input_file_entry = ttk.Entry(file_card, width=50)
    input_file_entry.grid(row=0, column=1, padx=10)
    ttk.Button(file_card, text="选择文件", command=lambda: (input_file_entry.delete(0, tk.END), input_file_entry.insert(0, filedialog.askopenfilename()))).grid(row=0, column=2)
    ttk.Label(file_card, text="保存位置:").grid(row=1, column=0, sticky="w", pady=5)
    output_folder_entry = ttk.Entry(file_card, width=50)
    output_folder_entry.grid(row=1, column=1, padx=10)
    ttk.Button(file_card, text="选择目录", command=lambda: (output_folder_entry.delete(0, tk.END), output_folder_entry.insert(0, filedialog.askdirectory()))).grid(row=1, column=2)
    ttk.Label(file_card, text="自定义文件名:").grid(row=2, column=0, sticky="w", pady=5)
    custom_filename_entry = ttk.Entry(file_card, width=50)
    custom_filename_entry.grid(row=2, column=1, padx=10, columnspan=2, sticky="w")
    dir_card = ttk.LabelFrame(main_frame, text=" 翻译语种 ", padding=15)
    dir_card.pack(fill="x", pady=10)
    translation_direction = tk.StringVar(value='ko2zh')
    lang_grid = ttk.Frame(dir_card)
    lang_grid.pack(fill="x")
    langs = [("韩 -> 中", 'ko2zh'), ("中 -> 韩", 'zh2ko'), ("韩 -> 越", 'ko2vi'), ("韩 -> 英", 'ko2en'),
             ("中 -> 英", 'zh2en'), ("英 -> 中", 'en2zh'), ("繁中 -> 英", 'zh_tw2en'), ("英 -> 繁中", 'en2zh_tw'),
             ("中 -> 日", 'zh2ja'), ("日 -> 中", 'ja2zh'), ("英 -> 韩", 'en2ko'), ("越 -> 中", 'vi2zh')]
    for i, (text, val) in enumerate(langs):
        row, col = i // 4, i % 4
        ttk.Radiobutton(lang_grid, text=text, variable=translation_direction, value=val).grid(row=row, column=col, padx=15, pady=5, sticky="w")
    opt_frame = tk.Frame(main_frame, bg="#f5f6fa")
    opt_frame.pack(fill="x", pady=10)
    append_translation = tk.BooleanVar(value=False)
    tk.Checkbutton(opt_frame, text="在原文下方保留翻译对照", variable=append_translation, bg="#f5f6fa", selectcolor="white").pack(side="left", padx=20)
    generate_corpus = tk.BooleanVar(value=False)
    tk.Checkbutton(opt_frame, text="同步生成语料库 (Corpus)", variable=generate_corpus, bg="#f5f6fa", selectcolor="white").pack(side="left", padx=20)
    translate_button = tk.Button(main_frame, text="🚀 开始长句翻译任务", command=start_translation, bg="#3498db", fg="white", font=("微软雅黑", 12, "bold"), relief="flat", cursor="hand2", pady=12)
    translate_button.pack(fill="x", pady=20)
    status_label = ttk.Label(main_frame, text="就绪：校准文件已自动挂载 (revision.md)", foreground="#7f8c8d")
    status_label.pack()

    revision_map = load_revision_dict("revision.md")
    root.mainloop()

if __name__ == "__main__":
    if TK_AVAILABLE and len(sys.argv) == 1:
        _run_gui()
    else:
        _run_cli()
