import warnings
# 静默 openpyxl 关于 DrawingML 支持不全的警告
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl.reader.drawings")

import sys
import argparse
try:
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk
    TK_AVAILABLE = True
except Exception:
    tk = None
    filedialog = None
    messagebox = None
    ttk = None
    TK_AVAILABLE = False
from pptx import Presentation
try:
    import pandas as pd
except Exception:
    pd = None
import os
import json
import hashlib
import random
import string
import requests
from openpyxl import load_workbook, Workbook
from openpyxl.styles import Font, Border, Alignment, PatternFill
from docx import Document
from docx.enum.text import WD_BREAK
from docx.text.paragraph import Paragraph
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.oxml.ns import nsmap as DOCX_NSMAP
from datetime import datetime
import re
import time
import threading
import shutil
import copy

# 请替换为你自己的 APP ID 和密钥（建议通过环境变量注入，避免提交到仓库）
APP_ID = os.getenv('BAIDU_APP_ID') or 'YOUR_APP_ID'
SECRET_KEY = os.getenv('BAIDU_SECRET_KEY') or 'YOUR_SECRET_KEY'
API_URL = 'https://api.fanyi.baidu.com/api/trans/vip/translate'

def configure_baidu(app_id=None, secret_key=None):
    global APP_ID, SECRET_KEY
    if app_id:
        APP_ID = str(app_id).strip()
    if secret_key:
        SECRET_KEY = str(secret_key).strip()

def _load_local_baidu_credentials():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    for filename in ("baidu_credentials.json", "baidu_credentials.local.json"):
        path = os.path.join(base_dir, filename)
        if not os.path.exists(path):
            continue
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f) or {}
            app_id = str(data.get("BAIDU_APP_ID") or "").strip()
            secret_key = str(data.get("BAIDU_SECRET_KEY") or "").strip()
            if app_id and secret_key:
                return app_id, secret_key
        except Exception:
            continue
    return None, None

if APP_ID == "YOUR_APP_ID" or SECRET_KEY == "YOUR_SECRET_KEY":
    _app_id, _secret_key = _load_local_baidu_credentials()
    if _app_id and _secret_key:
        APP_ID = _app_id
        SECRET_KEY = _secret_key

# ==========================================
# V2.12 (PPT 更新）
# ==========================================

# 全局变量
revision_map = {}
original_texts = []
translated_texts = []

class _ValueBox:
    def __init__(self, value=None):
        self._value = value
    def get(self):
        return self._value
    def set(self, value):
        self._value = value
# 扩大韩文匹配范围：包括预组合音节、基础辅音/元音、扩展字母等
HANGUL_RE = re.compile(r"[\uac00-\ud7af\u1100-\u11ff\u3130-\u318f\ua960-\ua97f\ud7b0-\ud7ff]")
# 扩大中文匹配范围：包括基本、扩展A、兼容汉字等
CHINESE_RE = re.compile(r"[\u4e00-\u9fff\u3400-\u4dbf\uf900-\ufaff]")
_translation_cache = {}  # 内存缓存：记录已翻译过的文本，避免重复请求 API

def should_skip_translation(text, to_lang):
    """
    智能判断是否跳过翻译：
    1. 如果目标是中文且原文已含中文，则跳过（满足用户：即使含部分韩文也不翻译中文）
    2. 如果原文和目标语种特征一致，则跳过
    """
    if not text or not str(text).strip():
        return True
    
    # 转换为字符串处理
    text_s = str(text).strip()
    
    # 核心逻辑：如果目标是中文（简体/繁体），且原文中已经包含中文字符，则跳过翻译直接返回原文
    if to_lang in ('zh', 'zh_tw') and CHINESE_RE.search(text_s):
        return True
        
    # 如果目标是韩文，且原文全是韩文（不含中文），也可以考虑跳过
    if to_lang == 'kor' and HANGUL_RE.search(text_s) and not CHINESE_RE.search(text_s):
        return True
        
    # 如果目标是英文，且原文不含中韩文，则跳过
    if to_lang == 'en' and not HANGUL_RE.search(text_s) and not CHINESE_RE.search(text_s):
        return True
        
    return False

FROM_LANG_MAP = {
    'zh2ko': 'zh', 'ko2zh': 'kor', 'ko2en': 'kor', 'zh2en': 'zh', 'en2zh': 'en',
    'zh_tw2en': 'zh', 'en2zh_tw': 'en', 'zh2ja': 'zh', 'ja2zh': 'ja',
    'en2ko': 'en', 'vi2zh': 'vie', 'zh2vi': 'zh', 'ko2vi': 'kor'
}
TO_LANG_MAP = {
    'zh2ko': 'kor', 'ko2zh': 'zh', 'ko2en': 'en', 'zh2en': 'en', 'en2zh': 'zh',
    'zh_tw2en': 'en', 'en2zh_tw': 'zh', 'zh2ja': 'ja', 'ja2zh': 'zh',
    'en2ko': 'kor', 'vi2zh': 'zh', 'zh2vi': 'vie', 'ko2vi': 'vie'
}
TARGET_FONT_BY_TO_LANG = {'zh': '微软雅黑', 'kor': 'Malgun Gothic'}
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def get_target_to_lang():
    try:
        direction = translation_direction.get()
    except Exception:
        return None
    return TO_LANG_MAP.get(direction)

def get_target_font_name():
    to_lang = get_target_to_lang()
    if not to_lang:
        return None
    return TARGET_FONT_BY_TO_LANG.get(to_lang)

def xpath_with_ns(element, expr):
    if element is None:
        return []
    try:
        return element.xpath(expr, namespaces=DOCX_NSMAP)
    except TypeError:
        return element.xpath(expr)

def set_docx_r_element_font(r_element, font_name):
    if r_element is None or not font_name:
        return
    rPr = r_element.get_or_add_rPr()
    rFonts = rPr.rFonts
    if rFonts is None:
        rFonts = OxmlElement('w:rFonts')
        rPr.append(rFonts)
    rFonts.set(qn('w:ascii'), font_name)
    rFonts.set(qn('w:hAnsi'), font_name)
    rFonts.set(qn('w:eastAsia'), font_name)
    rFonts.set(qn('w:cs'), font_name)
    for k in (qn('w:asciiTheme'), qn('w:hAnsiTheme'), qn('w:eastAsiaTheme'), qn('w:csTheme')):
        try:
            rFonts.attrib.pop(k, None)
        except Exception:
            pass

def set_docx_run_font(run, font_name):
    if not run or not font_name:
        return
    run.font.name = font_name
    set_docx_r_element_font(run._r, font_name)

def apply_run_format(source_run, target_run):
    """将 source_run 的所有格式（rPr）深度应用到 target_run"""
    if not source_run or not target_run:
        return
    source_rPr = source_run._element.rPr
    if source_rPr is not None:
        new_rPr = copy.deepcopy(source_rPr)
        target_r = target_run._element
        if target_r.rPr is not None:
            target_r.remove(target_r.rPr)
        target_r.insert(0, new_rPr)

def set_drawingml_rpr_element_font(a_rPr_element, font_name):
    if a_rPr_element is None or not font_name:
        return
    for local_name, tag in (('latin', 'a:latin'), ('ea', 'a:ea'), ('cs', 'a:cs')):
        el = a_rPr_element.find(f'{{{A_NS}}}{local_name}')
        if el is None:
            el = OxmlElement(tag)
            a_rPr_element.append(el)
        el.set('typeface', font_name)

def set_drawingml_r_element_font(a_r_element, font_name):
    if a_r_element is None or not font_name:
        return
    a_rPr = a_r_element.find(f'{{{A_NS}}}rPr')
    if a_rPr is None:
        a_rPr = OxmlElement('a:rPr')
        a_r_element.insert(0, a_rPr)
    set_drawingml_rpr_element_font(a_rPr, font_name)

def set_pptx_run_font(run, font_name):
    if not run or not font_name:
        return
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        r = run._r
        if hasattr(r, 'get_or_add_rPr'):
            rPr = r.get_or_add_rPr()
        else:
            rPr = r.find(f'{{{A_NS}}}rPr')
            if rPr is None:
                rPr = OxmlElement('a:rPr')
                r.insert(0, rPr)
        set_drawingml_rpr_element_font(rPr, font_name)
    except Exception:
        return

def set_pptx_paragraph_default_font(paragraph, font_name):
    if not paragraph or not font_name:
        return
    try:
        p = paragraph._p
        pPr = p.find(f'{{{A_NS}}}pPr')
        if pPr is None:
            pPr = OxmlElement('a:pPr')
            p.insert(0, pPr)
        defRPr = pPr.find(f'{{{A_NS}}}defRPr')
        if defRPr is None:
            defRPr = OxmlElement('a:defRPr')
            pPr.append(defRPr)
        set_drawingml_rpr_element_font(defRPr, font_name)
    except Exception:
        return

def load_revision_dict(file_path="revision.md", silent=False):
    """
    自动加载校准词典 (revision.md)
    """
    mapping = {}
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line or line.startswith("#"): continue
                    if not re.match(r'^[-*+]\s*', line): continue
                    content = re.sub(r'^[-*+]\s*', '', line)
                    if "格式" in content: continue
                    
                    if ":" in content:
                        parts = content.split(":", 1)
                        if len(parts) == 2:
                            err, corr = parts
                            if err.strip(): mapping[err.strip()] = corr.strip()
                    elif "：" in content:
                        parts = content.split("：", 1)
                        if len(parts) == 2:
                            err, corr = parts
                            if err.strip(): mapping[err.strip()] = corr.strip()
            if not silent:
                print(f"[系统] 自动加载校准文件成功，共 {len(mapping)} 条有效规则。")
        except Exception as e:
            print(f"[错误] 加载校准文件失败: {e}")
    return mapping

def apply_revisions(text):
    """
    在翻译结果上应用校准映射
    """
    if not text or not revision_map: return text
    sorted_keys = sorted(revision_map.keys(), key=len, reverse=True)
    result_text = text
    for err in sorted_keys:
        if err in result_text:
            result_text = result_text.replace(err, revision_map[err])
    return result_text

def baidu_translate(q, from_lang, to_lang):
    if not q or not q.strip(): return q
    
    # 1. 检查内存缓存
    cache_key = (from_lang, to_lang, str(q))
    cached = _translation_cache.get(cache_key)
    if cached is not None:
        return cached

    max_retries = 5
    retry_delay = 1.5
    retries = 0
    while retries < max_retries:
        try:
            salt = ''.join(random.choices(string.ascii_letters + string.digits, k=16))
            sign_str = APP_ID + q + salt + SECRET_KEY
            sign = hashlib.md5(sign_str.encode()).hexdigest()
            params = {
                'q': q, 'from': from_lang, 'to': to_lang,
                'appid': APP_ID, 'salt': salt, 'sign': sign
            }
            response = requests.get(API_URL, params=params, timeout=15)
            response.encoding = 'utf-8'
            result = response.json()
            
            if 'trans_result' in result:
                # 【核心修复】百度 API 返回的是一个列表，必须合并所有段落结果
                # 否则如果输入包含换行符，只会得到第一行的翻译
                dst_list = [item['dst'] for item in result['trans_result']]
                dst = "\n".join(dst_list)
                _translation_cache[cache_key] = dst  # 3. 写入缓存
                return dst
            
            # 处理百度 API 错误码
            error_code = str(result.get('error_code', ''))
            if error_code:
                error_msg = result.get('error_msg', 'Unknown Error')
                # 54003: 访问频率受限, 54005: 长查询请求频繁
                if error_code in ('54003', '54005'):
                    retries += 1
                    wait_time = retry_delay * (2 ** (retries - 1))
                    print(f"[警告] 百度 API 频率限制 ({error_code}: {error_msg})，第 {retries} 次重试，等待 {wait_time:.1f}s...")
                    time.sleep(wait_time)
                    continue
                else:
                    print(f"[错误] 百度 API 报错 ({error_code}: {error_msg})")
                    break # 其他错误（如 52003 账户欠费）不再重试
            
            return q
        except Exception as e:
            retries += 1
            if retries < max_retries:
                wait_time = retry_delay * (2 ** (retries - 1))
                print(f"[错误] 网络异常 ({type(e).__name__})，第 {retries} 次重试...")
                time.sleep(wait_time)
            else:
                print(f"[错误] 翻译请求最终失败: {e}")
    return q

def get_translation(text):
    if not text or not text.strip(): return text
    direction = translation_direction.get()
    from_lang = FROM_LANG_MAP.get(direction, 'auto')
    to_lang = TO_LANG_MAP.get(direction, 'auto')
    
    # --- [新增] 智能跳过逻辑 ---
    if should_skip_translation(text, to_lang):
        return text
    
    translated_text = baidu_translate(text, from_lang, to_lang)
    # 事后校准
    translated_text = apply_revisions(translated_text)
    return translated_text

def _ppt_normalize_linebreaks(text):
    if not text:
        return text
    return str(text).replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")

def _ppt_denormalize_linebreaks(text):
    if not text:
        return text
    return str(text).replace("\n", "\v")

def translate_ppt_paragraph(paragraph):
    full_text = paragraph.text
    if not full_text or not str(full_text).strip():
        return

    normalized = _ppt_normalize_linebreaks(full_text)
    translated = get_translation(normalized)

    original_texts.append(normalized)
    translated_texts.append(translated)

    if append_translation.get():
        result_norm = append_translation_to_original(normalized, translated)
    else:
        result_norm = translated

    result_text = _ppt_denormalize_linebreaks(result_norm)

    if paragraph.runs:
        first_run = paragraph.runs[0]
    else:
        first_run = paragraph.add_run()

    original_font_name = first_run.font.name
    original_size = first_run.font.size
    original_bold = first_run.font.bold
    original_italic = first_run.font.italic
    original_underline = first_run.font.underline

    first_run.text = result_text
    for r in paragraph.runs[1:]:
        r.text = ""

    target_font_name = get_target_font_name()
    if target_font_name:
        set_pptx_run_font(first_run, target_font_name)
        set_pptx_paragraph_default_font(paragraph, target_font_name)
    elif original_font_name:
        first_run.font.name = original_font_name
        first_run.font.size = original_size
        first_run.font.bold = original_bold
        first_run.font.italic = original_italic
        first_run.font.underline = original_underline

def append_translation_to_original(text, translated_text, cell=None):
    text = text.strip()
    translated_text = translated_text.strip()
    result = f"{text}\n{translated_text}" if text and translated_text else (text or translated_text)
    if cell:
        cell.alignment = Alignment(wrap_text=True, vertical='center')
        ws = cell.parent
        row_num = cell.row
        original_height = ws.row_dimensions[row_num].height
        ws.row_dimensions[row_num].height = (original_height * 2) if original_height else 30
    return result

# ==========================================
# 文件处理逻辑 (完全基于 V2.9 的稳定代码)
# ==========================================

def translate_shape_for_ppt(shape):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            translate_ppt_paragraph(paragraph)
    elif shape.shape_type == 6:  # 组合形状
        for sub_shape in shape.shapes:
            translate_shape_for_ppt(sub_shape)
    elif shape.has_table:  # 处理表格形状
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    translate_ppt_paragraph(paragraph)

def update_ui_status(msg):
    """线程安全地更新 UI 状态"""
    root.after(0, lambda: status_label.config(text=msg))

def translate_ppt(input_file, output_file):
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        msg = f"正在翻译 PPT: 第 {i}/{total_slides} 页..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        for shape in slide.shapes:
            translate_shape_for_ppt(shape)
    
    print(f"[系统] 正在保存 PPT 文件...")
    prs.save(output_file)
    save_to_corpus(original_texts, translated_texts)
    print(f"[完成] 保存修改后的 PPT 文件完成!")

def clean_sheet_name(name):
    if not name: return "Sheet"
    invalid_chars = r'[\\/?:*\[\](){}<>|"\']'
    clean_name = re.sub(invalid_chars, '', name)
    return clean_name[:31]

def translate_excel_xlsx(input_file, output_file):
    """处理.xlsx格式的Excel文件"""
    # 加载工作簿 (keep_vba=True 尝试保留宏和部分绘图元数据)
    wb = load_workbook(input_file, keep_vba=True)
    # 收集所有已存在的工作表名称（包括原始和已处理的）
    existing_sheet_names = set()
    total_sheets = len(wb.sheetnames)
    for i, sheet_name in enumerate(wb.sheetnames, 1):
        msg = f"正在翻译 Excel: 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        
        # 翻译工作表名称 (还原 V2.9 逻辑并增加翻译)
        translated_sheet_name = get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        
        # 确保工作表名称唯一
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)
        
        ws = wb[sheet_name]
        # 修改工作表名称
        ws.title = unique_sheet_name
        
        # 获取总行数用于进度（大致）
        total_rows = ws.max_row
        for row_idx, row in enumerate(ws.iter_rows(), 1):
            if row_idx % 10 == 0:
                row_msg = f"处理中: {unique_sheet_name} 第 {row_idx}/{total_rows} 行..."
                print(f"[进度] {row_msg}")
                update_ui_status(row_msg)
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    original_value = str(cell.value)
                    normalized_value = original_value.replace("\r\n", "\n").replace("\r", "\n")
                    if "\n" in normalized_value:
                        translated_lines = []
                        for line in normalized_value.split("\n"):
                            if line.strip():
                                translated_lines.append(get_translation(line))
                            else:
                                translated_lines.append(line)
                        translated_text = "\n".join(translated_lines)
                    else:
                        translated_text = get_translation(normalized_value)

                    original_texts.append(normalized_value)
                    translated_texts.append(translated_text)
                    
                    # 严格还原 V2.9 样式保护逻辑
                    # 保存原格式
                    original_font = Font(**cell.font.__dict__)
                    original_border = Border(**cell.border.__dict__)
                    original_alignment = Alignment(**cell.alignment.__dict__)
                    original_fill = PatternFill(**cell.fill.__dict__)
                    
                    if append_translation.get():
                        cell.value = append_translation_to_original(normalized_value, translated_text, cell)
                    else:
                        cell.value = translated_text
                    
                    # 恢复原格式
                    target_font_name = get_target_font_name()
                    if target_font_name:
                        cell.font = Font(
                            name=target_font_name,
                            size=original_font.size,
                            bold=original_font.bold,
                            italic=original_font.italic,
                            underline=original_font.underline,
                            color=original_font.color,
                        )
                    else:
                        cell.font = original_font
                    cell.border = original_border
                    cell.alignment = original_alignment
                    cell.fill = original_fill
                    # 设置单元格自动换行
                    cell.alignment = Alignment(wrap_text=True, vertical='center')
    # 保存修改后的工作簿
    wb.save(output_file)
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)

def translate_excel_xls(input_file, output_file):
    """处理.xls格式的Excel文件，使用pandas库，保存为.xlsx格式"""
    base_output = os.path.splitext(output_file)[0]
    output_file_xlsx = base_output + '.xlsx'
    target_font_name = get_target_font_name()
    # 读取所有工作表
    excel_file = pd.ExcelFile(input_file)
    # 创建一个新的ExcelWriter对象，使用默认引擎
    writer = pd.ExcelWriter(output_file_xlsx, engine='openpyxl')
    # 收集所有已存在的工作表名称（包括原始和已处理的）
    existing_sheet_names = set()
    total_sheets = len(excel_file.sheet_names)
    for i, sheet_name in enumerate(excel_file.sheet_names, 1):
        msg = f"正在翻译 Excel(.xls): 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        
        # 翻译工作表名称
        translated_sheet_name = get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        
        # 确保工作表名称唯一
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)
        
        # 读取工作表
        df = pd.read_excel(input_file, sheet_name=sheet_name)
        # 遍历所有单元格进行翻译
        total_rows = len(df.index)
        for idx_num, idx in enumerate(df.index, 1):
            if idx_num % 10 == 0:
                row_msg = f"处理中: {unique_sheet_name} 第 {idx_num}/{total_rows} 行..."
                print(f"[进度] {row_msg}")
                update_ui_status(row_msg)
            for col in df.columns:
                cell_value = df.at[idx, col]
                if pd.notna(cell_value) and isinstance(cell_value, str):
                    original_value = str(cell_value)
                    normalized_value = original_value.replace("\r\n", "\n").replace("\r", "\n")
                    if "\n" in normalized_value:
                        translated_lines = []
                        for line in normalized_value.split("\n"):
                            if line.strip():
                                translated_lines.append(get_translation(line))
                            else:
                                translated_lines.append(line)
                        translated_text = "\n".join(translated_lines)
                    else:
                        translated_text = get_translation(normalized_value)

                    original_texts.append(normalized_value)
                    translated_texts.append(translated_text)

                    df.at[idx, col] = append_translation_to_original(normalized_value, translated_text) if append_translation.get() else translated_text
        # 保存翻译后的工作表
        df.to_excel(writer, sheet_name=unique_sheet_name, index=False)
    # 保存工作簿
    writer.close()
    if target_font_name:
        try:
            wb = load_workbook(output_file_xlsx)
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            cell.font = Font(
                                name=target_font_name,
                                size=cell.font.size,
                                bold=cell.font.bold,
                                italic=cell.font.italic,
                                underline=cell.font.underline,
                                color=cell.font.color,
                            )
            wb.save(output_file_xlsx)
        except Exception:
            pass
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)
    return output_file_xlsx

def translate_word(input_file, output_file):
    doc = Document(input_file)
    target_font_name = get_target_font_name()
    # 核心：使用集合记录已处理的 XML 元素对象，防止重复翻译
    processed_elements = set()

    def set_style_font(style_element, font_name):
        if style_element is None or not font_name:
            return
        rPr = style_element.find(qn('w:rPr'))
        if rPr is None:
            rPr = OxmlElement('w:rPr')
            style_element.append(rPr)
        rFonts = rPr.find(qn('w:rFonts'))
        if rFonts is None:
            rFonts = OxmlElement('w:rFonts')
            rPr.append(rFonts)
        rFonts.set(qn('w:ascii'), font_name)
        rFonts.set(qn('w:hAnsi'), font_name)
        rFonts.set(qn('w:eastAsia'), font_name)
        rFonts.set(qn('w:cs'), font_name)
        for k in (qn('w:asciiTheme'), qn('w:hAnsiTheme'), qn('w:eastAsiaTheme'), qn('w:csTheme')):
            try:
                rFonts.attrib.pop(k, None)
            except Exception:
                pass

    def translate_word_paragraph(paragraph):
        """通用段落翻译逻辑：支持对照模式、换行符、字体设置、图片保留"""
        p_id = id(paragraph._element)
        if p_id in processed_elements:
            return
        processed_elements.add(p_id)

        full_text = paragraph.text
        if not full_text or not str(full_text).strip():
            return

        # --- 捕捉原句的模板 Run (用于对照模式下的格式克隆) ---
        template_run = paragraph.runs[0] if paragraph.runs else None
        
        # 使用 splitlines() 更加稳健地处理 \r, \n, \v 等换行符
        lines = [line.strip() for line in full_text.splitlines() if line.strip()]
        if not lines:
            return

        print(f"[处理] Word 段落 ({len(lines)} 行): {lines[0][:30]}...")

        if append_translation.get():
            # 【对照模式优化：行对行对照】
            translated_lines = []
            for line in lines:
                t = get_translation(line)
                translated_lines.append(t)
            
            # 记录语料库
            for orig, trans in zip(lines, translated_lines):
                original_texts.append(orig)
                translated_texts.append(trans)
            
            # 清空段落所有内容
            p_el = paragraph._element
            for r in p_el.xpath('./w:r'):
                p_el.remove(r)
            
            # 重新填充
            for i, (orig, trans) in enumerate(zip(lines, translated_lines)):
                # 1. 写入原文 (应用模板格式)
                run_orig = paragraph.add_run(orig)
                if template_run:
                    apply_run_format(template_run, run_orig)
                
                run_orig.add_break(WD_BREAK.LINE)
                
                # 2. 写入译文 (应用模板格式，并覆盖目标字体)
                run_trans = paragraph.add_run(trans)
                if template_run:
                    apply_run_format(template_run, run_trans)
                
                if target_font_name:
                    set_docx_run_font(run_trans, target_font_name)
                
                # 组间换行
                if i < len(lines) - 1:
                    run_trans.add_break(WD_BREAK.LINE)
        else:
            # 【替换模式】
            normalized = "\n".join(lines)
            t = get_translation(normalized)
            original_texts.append(normalized)
            translated_texts.append(t)
            
            if paragraph.runs:
                first_run = paragraph.runs[0]
                first_run.text = t
                for r in paragraph.runs[1:]:
                    if not bool(xpath_with_ns(r.element, './/w:drawing')):
                        r.text = ""
                if target_font_name:
                    set_docx_run_font(first_run, target_font_name)
            else:
                new_run = paragraph.add_run(t)
                if target_font_name:
                    set_docx_run_font(new_run, target_font_name)

    def process_xml_container(container_el):
        """深度遍历 XML 容器中的所有段落和文本节点，确保顺序一致"""
        # 1. 查找所有标准段落 <w:p>
        # 这涵盖了：正文、表格单元格、文本框中的所有段落
        for p_el in container_el.xpath('.//w:p'):
            try:
                p_obj = Paragraph(p_el, doc)
                translate_word_paragraph(p_obj)
            except Exception as e:
                print(f"[警告] 处理段落 XML 失败: {e}")

        # 2. 查找所有图形文字段落 <a:p> (DrawingML)
        # 针对图片内部、形状内部的特殊文字结构
        for p_el in container_el.xpath('.//a:p'):
            if p_el in processed_elements:
                continue
            processed_elements.add(p_el)
            
            t_nodes = p_el.xpath('.//a:t')
            if not t_nodes: continue
            
            combined_text = "".join([node.text for node in t_nodes if node.text])
            if not combined_text.strip(): continue
            
            t = get_translation(combined_text)
            original_texts.append(combined_text)
            translated_texts.append(t)
            
            if append_translation.get():
                # DrawingML 换行修复
                t_nodes[0].text = combined_text
                for n in t_nodes[1:]: n.text = ""
                r_nodes = p_el.xpath('.//a:r')
                if r_nodes:
                    last_r = r_nodes[0]
                    br = OxmlElement('a:br')
                    last_r.addnext(br)
                    new_r = OxmlElement('a:r')
                    new_t = OxmlElement('a:t')
                    new_t.text = t
                    new_r.append(new_t)
                    br.addnext(new_r)
                    if target_font_name:
                        set_drawingml_r_element_font(new_r, target_font_name)
            else:
                t_nodes[0].text = t
                for n in t_nodes[1:]: n.text = ""

    # --- 执行统一遍历 ---
    
    # A. 处理正文内容 (含表格、文本框、图形)
    print(f"[系统] 正在按顺序扫描 Word 正文...")
    process_xml_container(doc.element.body)

    # B. 处理所有页眉和页脚
    for i, section in enumerate(doc.sections, 1):
        for header in [section.header, section.first_page_header, section.even_page_header]:
            if header: process_xml_container(header._element)
        for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
            if footer: process_xml_container(footer._element)

    # C. 全局字体统一加固 (仅针对非对照模式，或者作为样式兜底)
    if target_font_name:
        try:
            # 1. 优先设置样式表字体
            for s in doc.styles:
                try:
                    if getattr(s, "font", None) is not None:
                        s.font.name = target_font_name
                except Exception: pass
                try:
                    set_style_font(getattr(s, "_element", None), target_font_name)
                except Exception: pass
            
            # 2. 注意：移除对所有 w:r 的强制刷字体，防止覆盖对照模式下的原文原字体
            # 如果是非对照模式，可以保留这个全局刷
            if not append_translation.get():
                for r in doc.element.xpath('.//w:r'):
                    set_docx_r_element_font(r, target_font_name)
            
            # 图形文字 (DrawingML) 通常需要强制刷，因为它们往往没有复杂的混合样式
            for r in doc.element.xpath('.//a:r'):
                set_drawingml_r_element_font(r, target_font_name)
        except Exception: pass

    doc.save(output_file)
    save_to_corpus(original_texts, translated_texts)

# --- 线程控制 ---

def _extract_sample_text(input_file):
    """
    极简采样：只提取文件最开始的一点内容
    """
    try:
        sample_text = ""
        ext = os.path.splitext(input_file)[1].lower()
        if ext == '.docx':
            doc = Document(input_file)
            for p in doc.paragraphs:
                if p.text.strip():
                    sample_text = p.text
                    break
        elif ext == '.xlsx':
            wb = load_workbook(input_file, read_only=True, data_only=True)
            ws = wb.active
            for row in ws.iter_rows(max_row=10):
                for cell in row:
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        sample_text = str(cell.value)
                        break
                if sample_text: break
        elif ext == '.ppt' or ext == '.pptx':
            prs = Presentation(input_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        sample_text = shape.text_frame.text
                        break
                if sample_text: break
        return sample_text
    except Exception:
        return ""

def _judge_direction_warning(langs, direction):
    """
    极简逻辑：只看原语言是否一致
    """
    if not langs: return None
    
    from_lang = FROM_LANG_MAP.get(direction)
    from_lang_names = {'zh': '中文', 'kor': '韩文', 'ja': '日文', 'en': '英文', 'vi': '越文'}
    source_name = from_lang_names.get(from_lang, from_lang)

    # 逻辑：如果不包含所选的源语言，直接拦截
    if from_lang == 'zh' and not langs.get("has_chinese"):
        return f"文件开头未检测到 [{source_name}]"
    elif from_lang == 'kor' and not langs.get("has_korean"):
        return f"文件开头未检测到 [{source_name}]"
    elif from_lang == 'ja' and not langs.get("has_japanese"):
        return f"文件开头未检测到 [{source_name}]"
    elif from_lang == 'en' and not langs.get("has_english"):
        return f"文件开头未检测到 [{source_name}]"

    return None

def check_direction_mismatch(input_file, direction=None):
    """
    预检：如果翻译方向和文件内容明显不匹配，返回错误原因
    """
    try:
        if direction is None:
            try:
                direction = translation_direction.get()
            except Exception:
                return None
        
        sample_text = _extract_sample_text(input_file)
        if not sample_text.strip():
            return None

        langs = {
            "has_korean": bool(HANGUL_RE.search(sample_text)),
            "has_chinese": bool(CHINESE_RE.search(sample_text)),
            "has_japanese": bool(re.search(r'[\u3040-\u30ff]', sample_text)),
            "has_english": bool(re.search(r'[a-zA-Z]{3,}', sample_text))
        }

        return _judge_direction_warning(langs, direction)
    except Exception as e:
        print(f"[预检] 预检过程出错: {e}")
        return None

def start_translation():
    input_file = input_file_entry.get()
    output_folder = output_folder_entry.get()
    if not (input_file and output_folder):
        messagebox.showwarning("提示", "请完整选择输入文件和输出目录")
        return
    
    # --- [新增] 翻译方向预检 ---
    warning_reason = check_direction_mismatch(input_file)
    if warning_reason:
        # 构造易读的方向标签
        dir_labels = {
            'ko2zh': '韩 -> 中', 'zh2ko': '中 -> 韩', 'ko2vi': '韩 -> 越', 'ko2en': '韩 -> 英',
            'zh2en': '中 -> 英', 'en2zh': '英 -> 中', 'zh_tw2en': '繁中 -> 英', 'en2zh_tw': '英 -> 繁中',
            'zh2ja': '中 -> 日', 'ja2zh': '日 -> 中', 'en2ko': '英 -> 韩', 'vi2zh': '越 -> 中'
        }
        current_dir_label = dir_labels.get(translation_direction.get(), translation_direction.get())
        if not messagebox.askyesno("预警", f"检测到当前翻译方向为 [{current_dir_label}]，但文件内容似乎是 {warning_reason}（既非源语言也非目标语言）。\n\n是否继续执行翻译？"):
            return
    
    translate_button.config(state=tk.DISABLED, text="🚀 正在翻译，请稍候...")
    status_label.config(text="任务已启动，请查看终端进度...", foreground="#2980b9")
    thread = threading.Thread(target=run_translation_task, args=(input_file, output_folder))
    thread.daemon = True
    thread.start()

def run_translation_task(input_file, output_folder):
    start_time = time.time()
    try:
        global revision_map
        file_ext = os.path.splitext(input_file)[1].lower()
        custom_name = custom_filename_entry.get().strip()
        output_file = os.path.join(output_folder, (custom_name if custom_name else f"translated_v2.12_{os.path.basename(input_file).split('.')[0]}") + file_ext)
        
        # --- [V2.12 增强功能：物理克隆以保留图片和绘图] ---
        # 即使 V2.9 也不支持 Excel 绘图保留，这里通过物理复制尝试最大化兼容性
        try:
            if os.path.exists(output_file):
                os.remove(output_file) # 尝试删除旧文件，如果被占用会在这里报错
            shutil.copy2(input_file, output_file)
        except PermissionError:
            raise Exception(f"目标文件已被占用，请先关闭 Excel/Word/PPT: {os.path.basename(output_file)}")
        
        original_texts.clear()
        translated_texts.clear()
        revision_map = load_revision_dict("revision.md", silent=True)
        print(f"\n[系统] 开始翻译任务: {os.path.basename(input_file)}")
        print(f"[系统] 校准规则加载: {len(revision_map)} 条")
        
        # 统一使用 output_file 作为操作对象，实现“原地翻译”
        if file_ext in ['.ppt', '.pptx']: translate_ppt(output_file, output_file)
        elif file_ext == '.xlsx': translate_excel_xlsx(output_file, output_file)
        elif file_ext == '.xls': 
            # .xls 比较特殊，必须另存为 .xlsx
            output_file = translate_excel_xls(input_file, output_file)
        elif file_ext == '.docx': translate_word(output_file, output_file)
        
        end_time = time.time()
        duration_minutes = (end_time - start_time) / 60
        print(f"[完成] 文件已保存至: {output_file}")
        print(f"[统计] 翻译总耗时: {duration_minutes:.2f} 分钟\n")
        root.after(0, lambda: translation_done_callback(output_file, duration_minutes))
    except Exception as e:
        err_msg = str(e)
        print(f"[错误] 详情: {err_msg}")
        root.after(0, lambda: translation_failed_callback(err_msg))

def translation_done_callback(output_file, duration_minutes):
    translate_button.config(state=tk.NORMAL, text="🚀 开始长句翻译任务")
    status_label.config(text=f"翻译任务已圆满完成！(用时 {duration_minutes:.2f} 分钟)", foreground="#27ae60")
    messagebox.showinfo("成功", f"翻译完成！\n用时：{duration_minutes:.2f} 分钟\n文件保存至：{output_file}")

def translation_failed_callback(error_msg):
    translate_button.config(state=tk.NORMAL, text="🚀 开始长句翻译任务")
    status_label.config(text=f"翻译过程出错", foreground="#e74c3c")
    messagebox.showerror("错误", f"发生异常：{error_msg}")

def save_to_corpus(orig, trans):
    if generate_corpus.get() and orig:
        if not os.path.exists('Corpus'): os.makedirs('Corpus')
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        try:
            direction = translation_direction.get()
        except Exception:
            direction = "unknown"
        corpus_file = f'Corpus/Corpus_v2.12_{direction}_{timestamp}.xlsx'
        to_lang = TO_LANG_MAP.get(direction)
        seen = set()
        filtered_orig = []
        filtered_trans = []
        for o, t in zip(orig, trans):
            if o is None or t is None:
                continue
            o = str(o).replace("\r\n", "\n").replace("\r", "\n").strip()
            t = str(t).replace("\r\n", "\n").replace("\r", "\n").strip()
            if not o or not t:
                continue

            # 3. 原文=译文 不记录
            if o == t:
                continue

            # 检查语言特征
            has_korean = bool(HANGUL_RE.search(o))
            has_chinese = bool(CHINESE_RE.search(o))

            # 1. 核心准则：只要含有韩文就要记录 (满足 "原文只要含有韩文就要记录")
            if has_korean:
                pass 
            else:
                # 如果不含韩文：
                # A. 如果包含中文，且目标是中文 (针对韩翻中方向) -> 不记录 (即 "纯中文不录")
                # 注意：如果包含中文但目标不是中文 (如中翻韩、中翻英)，则记录
                if has_chinese and to_lang in ("zh", "zh_tw"):
                    continue
                
                # B. 如果既不含韩文也不含中文 -> 说明只有数字、英文、符号的组合 -> 不记录
                # 这完全符合 "数字/英文/符号 只有全部是的时候才能不记录"
                if not has_chinese:
                    continue

            # 只要包含韩文，或者满足上述条件的非韩文行，都进入去重检查
            key = o
            if key in seen:
                continue
            seen.add(key)
            filtered_orig.append(o)
            filtered_trans.append(t)

        if not filtered_orig:
            return

        pd.DataFrame(
            {'序号': range(1, len(filtered_orig) + 1), '翻译前': filtered_orig, '翻译后': filtered_trans}
        ).to_excel(corpus_file, index=False)

def translate_file(input_file, output_dir=None, name="", direction="ko2zh", append=False, corpus=False, revision_file="revision.md", app_id=None, secret_key=None):
    global translation_direction, append_translation, generate_corpus, revision_map

    input_file = os.path.abspath(input_file)
    if not os.path.exists(input_file):
        raise FileNotFoundError(input_file)

    configure_baidu(app_id=app_id, secret_key=secret_key)

    output_dir = os.path.abspath(output_dir) if output_dir else os.path.dirname(input_file)
    os.makedirs(output_dir, exist_ok=True)

    translation_direction = _ValueBox(direction)
    append_translation = _ValueBox(bool(append))
    generate_corpus = _ValueBox(bool(corpus))

    original_texts.clear()
    translated_texts.clear()
    revision_map = load_revision_dict(revision_file, silent=True)

    file_ext = os.path.splitext(input_file)[1].lower()
    base_name = name.strip() if name and name.strip() else f"translated_{os.path.basename(input_file).split('.')[0]}"
    output_file = os.path.join(output_dir, base_name + file_ext)

    if os.path.exists(output_file):
        os.remove(output_file)
    shutil.copy2(input_file, output_file)

    if file_ext in ['.ppt', '.pptx']:
        translate_ppt(output_file, output_file)
    elif file_ext == '.xlsx':
        translate_excel_xlsx(output_file, output_file)
    elif file_ext == '.xls':
        output_file = translate_excel_xls(input_file, output_file)
    elif file_ext == '.docx':
        translate_word(output_file, output_file)
    else:
        raise ValueError(f"不支持的文件类型: {file_ext}")

    return output_file

# --- UI 布局 ---
if __name__ == "__main__" and TK_AVAILABLE:
    root = tk.Tk()
    root.title("百度长语句翻译工具 V2.12 (PPT 更新）")
    root.geometry("700x850")
    root.configure(bg="#f5f6fa")
    style = ttk.Style()
    style.theme_use('clam')
    style.configure("TFrame", background="#f5f6fa")
    style.configure("TLabel", background="#f5f6fa", font=("微软雅黑", 10))
    style.configure("Header.TLabel", font=("微软雅黑", 14, "bold"), foreground="#2f3640")
    main_frame = ttk.Frame(root, padding="20")
    main_frame.pack(fill="both", expand=True)
    ttk.Label(main_frame, text="百度长语句翻译 & 自动校准系统", style="Header.TLabel").pack(pady=(0, 20))
    file_card = ttk.LabelFrame(main_frame, text=" 文件设置 ", padding=15)
    file_card.pack(fill="x", pady=10)
    ttk.Label(file_card, text="待翻译文件:").grid(row=0, column=0, sticky="w", pady=5)
    input_file_entry = ttk.Entry(file_card, width=50)
    input_file_entry.grid(row=0, column=1, padx=10)
    ttk.Button(file_card, text="选择文件", command=lambda: (input_file_entry.delete(0, tk.END), input_file_entry.insert(0, filedialog.askopenfilename()))).grid(row=0, column=2)
    ttk.Label(file_card, text="保存位置:").grid(row=1, column=0, sticky="w", pady=5)
    output_folder_entry = ttk.Entry(file_card, width=50)
    output_folder_entry.grid(row=1, column=1, padx=10)
    ttk.Button(file_card, text="选择目录", command=lambda: (output_folder_entry.delete(0, tk.END), output_folder_entry.insert(0, filedialog.askdirectory()))).grid(row=1, column=2)
    ttk.Label(file_card, text="自定义文件名:").grid(row=2, column=0, sticky="w", pady=5)
    custom_filename_entry = ttk.Entry(file_card, width=50)
    custom_filename_entry.grid(row=2, column=1, padx=10, columnspan=2, sticky="w")
    dir_card = ttk.LabelFrame(main_frame, text=" 翻译语种 ", padding=15)
    dir_card.pack(fill="x", pady=10)
    translation_direction = tk.StringVar(value='ko2zh')
    lang_grid = ttk.Frame(dir_card)
    lang_grid.pack(fill="x")
    langs = [("韩 -> 中", 'ko2zh'), ("中 -> 韩", 'zh2ko'), ("韩 -> 越", 'ko2vi'), ("韩 -> 英", 'ko2en'),
             ("中 -> 英", 'zh2en'), ("英 -> 中", 'en2zh'), ("繁中 -> 英", 'zh_tw2en'), ("英 -> 繁中", 'en2zh_tw'),
             ("中 -> 日", 'zh2ja'), ("日 -> 中", 'ja2zh'), ("英 -> 韩", 'en2ko'), ("越 -> 中", 'vi2zh')]
    for i, (text, val) in enumerate(langs):
        row, col = i // 4, i % 4
        ttk.Radiobutton(lang_grid, text=text, variable=translation_direction, value=val).grid(row=row, column=col, padx=15, pady=5, sticky="w")
    opt_frame = tk.Frame(main_frame, bg="#f5f6fa")
    opt_frame.pack(fill="x", pady=10)
    append_translation = tk.BooleanVar(value=False)
    tk.Checkbutton(opt_frame, text="在原文下方保留翻译对照", variable=append_translation, bg="#f5f6fa", selectcolor="white").pack(side="left", padx=20)
    generate_corpus = tk.BooleanVar(value=False)
    tk.Checkbutton(opt_frame, text="同步生成语料库 (Corpus)", variable=generate_corpus, bg="#f5f6fa", selectcolor="white").pack(side="left", padx=20)
    translate_button = tk.Button(main_frame, text="🚀 开始长句翻译任务", command=start_translation, bg="#3498db", fg="white", font=("微软雅黑", 12, "bold"), relief="flat", cursor="hand2", pady=12)
    translate_button.pack(fill="x", pady=20)
    status_label = ttk.Label(main_frame, text="就绪：校准文件已自动挂载 (revision.md)", foreground="#7f8c8d")
    status_label.pack()

    # 初始化加载一次即可
    revision_map = load_revision_dict("revision.md")

    root.mainloop()
